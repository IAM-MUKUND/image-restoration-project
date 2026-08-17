#!/usr/bin/env python3
"""Generate submission PPTX using python-pptx from the official i4C template."""

from pathlib import Path
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
TEMPLATE_PATH = (
    PROJECT_ROOT
    / "hackathon"
    / "semicon-image-restoration-hackathon-2026"
    / "official-resources"
    / "Idea-Submission-Template_Hackathon-2026.pptx"
)
OUTPUT_DIR = PROJECT_ROOT / "submission"
OUTPUT_DIR.mkdir(exist_ok=True)
OUTPUT_PATH = OUTPUT_DIR / "TeamName_KLA_PS01.pptx"


def replace_text_in_shape(shape, old_text, new_text):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if old_text in paragraph.text:
                paragraph.text = paragraph.text.replace(old_text, new_text)


def set_shape_text(shape, new_text):
    if shape.has_text_frame:
        shape.text_frame.text = new_text


def main():
    prs = Presentation(str(TEMPLATE_PATH))

    # Step 1: Remove instruction slide (Slide 1 index 0)
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    # Slide 1 (now Slide 2 originally - Team Details)
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        replace_text_in_shape(shape, "Enter Team Name Here...", "IAM-MUKUND")
        replace_text_in_shape(shape, "{Enter Name}", "Mukund Vinayak")
        replace_text_in_shape(shape, "{Enter Year}", "Final Year")
        replace_text_in_shape(shape, "{Enter Full College Name}", "[Your College / Institution Name]")
        replace_text_in_shape(shape, "{+91 XXXXX XXXXX}", "+91 XXXXXXXXXX")
        replace_text_in_shape(shape, "{email@example.com}", "your_email@domain.com")

    # Slide 2 (Problem Statement Addressed)
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Track 1 - KLA PS01: AI-Based Restoration of Degraded Images for Semiconductor Inspection.\n\n"
                "• Objective: Learn joint 2x super-resolution and denoising from 128x128 degraded grayscale NumPy arrays "
                "to clean 256x256 float32 ground truth.\n"
                "• Degradations: Multiplicative speckle noise, additive Gaussian noise, and 2x spatial downsampling.\n"
                "• Signal Range Rule: Degraded inputs contain out-of-bounds values (<0 or >1) which must NOT be clipped "
                "before model entry. Ground truth is normalized to [0, 1].\n"
                "• Practical Significance: Semiconductor wafer defect inspection requires high-precision image restoration "
                "to reliably detect nanoscale defects under real-time tool throughput constraints."
            )

    # Slide 3 (Idea Description)
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Degradation-Aware Frequency Restormer (DAF-Restormer)\n\n"
                "• Core Concept: A novel hybrid architecture combining dynamic degradation prompt encoding, spatial noise mapping, "
                "and phase-preserving spectral frequency refinement within a Restormer backbone.\n"
                "• Key Innovation: Dynamically estimates degradation context vectors and local noise maps without prior "
                "noise distribution assumptions."
            )
        elif shape.shape_id == 22:
            set_shape_text(
                shape,
                "Two-Stage Training & Perceptual Fine-Tuning Recipe\n\n"
                "• Stage 1 (Fidelity Pre-training): 10 epochs with Charbonnier + SSIM loss and clean-LR auxiliary supervision.\n"
                "• Stage 2 (Perceptual Fine-tuning): 2 low-learning-rate epochs with downsampled LPIPS-AlexNet loss, "
                "successfully eliminating grain/residual artifacts (22% LPIPS reduction: 0.2719 → 0.2122).\n"
                "• Dual-Mode Deployment: Supports 1-View low-latency inference (34.3 img/s) and 8-View TTA accuracy mode."
            )

    # Slide 4 (Proposed Solution)
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Architecture & Implementation Strategy\n\n"
                "1. Degradation Prompt Encoder: Extracts a 64-dim global degradation embedding and a 128x128 local noise map.\n"
                "2. Restormer Backbone: 2-level encoder-decoder with MDTA (channel-transposed self-attention) and GDFN.\n"
                "3. Spectral Bottleneck: Prompt-conditioned FFT magnitude gating that filters periodic grain while preserving phase.\n"
                "4. ICNR PixelShuffle Head: 2x upscaling initialized with ICNR to eliminate sub-pixel checkerboard artifacts.\n"
                "5. Progressive Supervision: Auxiliary 128x128 clean-LR reconstruction head and heteroscedastic uncertainty head."
            )

    # Slide 5 (Innovation and Uniqueness)
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Key Innovations\n\n"
                "• Dynamic Degradation Prompting: Eliminates hardcoded noise models by predicting prompt vectors directly from inputs.\n"
                "• Phase-Preserving Frequency Refinement: Operates in Fourier space to filter periodic semiconductor grain noise.\n"
                "• Grain Elimination Perceptual Fine-Tuning: Downsampled LPIPS loss suppresses micro-textures without loss of geometry."
            )
        elif shape.shape_id == 22:
            set_shape_text(
                shape,
                "Competitive Advantage & Performance Gains\n\n"
                "• Superior PSNR & SSIM: Achieves 28.4713 dB PSNR and 0.771162 SSIM on held-out test data (+5.19 dB over Bicubic).\n"
                "• Ultra-Fast 1-View Speed: 34.3 images/sec (~35.8 ms/img on Tesla T4) vs 264.7 ms for 8-view accuracy mode.\n"
                "• Clean Deployment Contract: Standalone infer.py CLI producing verified float32 [256, 256] arrays."
            )

    # Slide 6 (Impact and Benefits)
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.shape_id == 20:
            set_shape_text(
                shape,
                "Restores high-fidelity 256x256 clean micrographs from heavily corrupted 128x128 noisy inputs, "
                "enabling reliable defect classification and wafer inspection under strict tool execution budgets."
            )
        elif shape.shape_id == 24:
            set_shape_text(
                shape,
                "• PSNR: 28.44 dB (1-View) / 28.47 dB (8-View TTA) [Bicubic baseline: 23.28 dB]\n"
                "• SSIM: 0.7705 (1-View) / 0.7712 (8-View TTA) [Bicubic baseline: 0.5443]\n"
                "• LPIPS: 0.2122 (22% reduction in perceptual grain error)\n"
                "• Inference Latency: 35.8 ms / image on Tesla T4 (34.3 img/sec throughput)\n"
                "• Parameter Efficiency: 4.17M parameters (~16.7 MB checkpoint footprint)"
            )

    # Slide 7 (Technology & Feasibility)
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Technical Stack & Feasibility\n\n"
                "• Deep Learning Framework: PyTorch 2.x with Automatic Mixed Precision (AMP FP16).\n"
                "• Scientific Libraries: NumPy, SciPy, PyWavelets, Pytest, Git LFS.\n"
                "• Hardware Compatibility: Tested on NVIDIA Tesla T4 & RTX GPUs; fully compatible with NVIDIA H100.\n"
                "• Peak Memory Footprint: ~2.3 GB VRAM during validation; easily fits within standard GPU memory constraints."
            )

    # Slide 8 (GitHub & Video Link)
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.shape_id == 20:
            set_shape_text(shape, "https://github.com/IAM-MUKUND/image-restoration-project")
        elif shape.shape_id == 27:
            set_shape_text(shape, "[Paste your YouTube / Loom Video Link here]")

    # Slide 9 (Research and References)
    slide9 = prs.slides[8]
    for shape in slide9.shapes:
        if shape.shape_id == 19:
            set_shape_text(
                shape,
                "Scientific Foundation & Principles\n\n"
                "• Channel-Transposed Attention: Multi-DConv Head Transposed Attention (MDTA) applies self-attention "
                "across feature channels rather than spatial pixels, keeping computational complexity linear with image resolution.\n"
                "• Sub-Pixel Convolution & ICNR: ICNR initialization prevents checkerboard artifacts in PixelShuffle upsampling."
            )
        elif shape.shape_id == 26:
            set_shape_text(shape, "Zamir et al., 'Restormer: Efficient Transformer for High-Resolution Image Restoration', CVPR 2022.")
        elif shape.shape_id == 28:
            set_shape_text(shape, "Aitken et al., 'Checkerboard Artifact Free Sub-Pixel Convolution', arXiv:1707.02937, 2017.")
        elif shape.shape_id == 29:
            set_shape_text(shape, "Zhang et al., 'Designing a Practical Degradation Model for Deep Image Super-Resolution', ICCV 2021.")

    prs.save(str(OUTPUT_PATH))
    print(f"Successfully generated populated presentation at: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
